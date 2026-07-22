"""
Сборка редактируемой PPTX-презентации из текущего состояния дашборда —
как ручной операционный отчёт клиники, но который можно доправить руками
в PowerPoint/Google Slides: титульный слайд, сводка с метриками и тезисами,
затем по 1-2 графика на слайд с заголовком-выводом.

Использует kaleido (рендер plotly-графиков в PNG) и python-pptx (сборка PPTX).
"""

import io
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

CREAM = RGBColor(0xFB, 0xF6, 0xF0)
BRAND_DEEP = RGBColor(0x21, 0x37, 0x4A)
GRAY_TEXT = RGBColor(0x3C, 0x3C, 0x3C)
RED_TEXT = RGBColor(0x78, 0x3C, 0x3C)
GREEN_TEXT = RGBColor(0x3C, 0x6E, 0x3C)

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5


def _fig_to_png(fig, path: str, width=1000, height=560, scale=2):
    fig.write_image(path, width=width, height=height, scale=scale)


def _set_cream_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = CREAM


def _add_logo(slide, logo_path: str | None):
    if logo_path and Path(logo_path).exists():
        slide.shapes.add_picture(logo_path, Inches(0.3), Inches(0.25), height=Inches(0.45))


def _add_textbox(slide, left, top, width, height, text, size, color, bold=False, italic=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return box


def build_pptx(logo_path: str, selected_period: str, data_range: str, metrics: list, narrative: list, sections: list, lang: str = "ru") -> bytes:
    """
    metrics: [(label, value_str, delta_str_or_None), ...]
    narrative: [(heading, text), ...]
    sections: [{"tab": str, "heading": str, "figs": [go.Figure, ...]}, ...]
    lang: "ru" или "en" — переключает статичные подписи слайдов (тезисы/заголовки
    секций уже приходят на нужном языке от вызывающего кода в app.py).
    """
    def _t(ru: str, en: str) -> str:
        return en if lang == "en" else ru

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank_layout = prs.slide_layouts[6]

    with tempfile.TemporaryDirectory() as tmp:
        tmp_path = Path(tmp)

        # ── Титульный слайд ──
        slide = prs.slides.add_slide(blank_layout)
        _set_cream_background(slide)
        _add_logo(slide, logo_path)
        _add_textbox(slide, 0, 2.8, SLIDE_W_IN, 1, _t("Операционный отчёт клиники", "Clinic operational report"), 34, BRAND_DEEP, bold=True, align=PP_ALIGN.CENTER)
        _add_textbox(slide, 0, 3.7, SLIDE_W_IN, 0.6, _t(f"Период: {selected_period}", f"Period: {selected_period}"), 18, GRAY_TEXT, align=PP_ALIGN.CENTER)
        _add_textbox(slide, 0, 4.2, SLIDE_W_IN, 0.5, _t(f"Данные: {data_range}", f"Data: {data_range}"), 11, GRAY_TEXT, italic=True, align=PP_ALIGN.CENTER)

        # ── Сводка: метрики + тезисы ──
        slide = prs.slides.add_slide(blank_layout)
        _set_cream_background(slide)
        _add_logo(slide, logo_path)
        _add_textbox(slide, 0.4, 0.5, 10, 0.6, _t(f"Главное за период: {selected_period}", f"Highlights for {selected_period}"), 22, BRAND_DEEP, bold=True)

        col_w = (SLIDE_W_IN - 0.8) / max(len(metrics), 1)
        for i, (label, value, delta) in enumerate(metrics):
            x = 0.4 + i * col_w
            _add_textbox(slide, x, 1.3, col_w - 0.1, 0.5, label, 10, GRAY_TEXT)
            _add_textbox(slide, x, 1.7, col_w - 0.1, 0.5, value, 16, BRAND_DEEP, bold=True)
            if delta:
                color = RED_TEXT if delta.startswith("-") else GREEN_TEXT
                _add_textbox(slide, x, 2.15, col_w - 0.1, 0.4, delta, 10, color)

        _add_textbox(slide, 0.4, 2.8, 10, 0.5, _t("Главные тезисы периода", "Key takeaways for the period"), 16, BRAND_DEEP, bold=True)
        y = 3.3
        for heading, text in narrative:
            _add_textbox(slide, 0.4, y, 12.5, 0.3, heading, 12, BRAND_DEEP, bold=True)
            y += 0.32
            box = _add_textbox(slide, 0.4, y, 12.5, 0.5, text, 11, GRAY_TEXT)
            y += 0.42

        # ── Секции с графиками ──
        current_tab = None
        for i, section in enumerate(sections):
            slide = prs.slides.add_slide(blank_layout)
            _set_cream_background(slide)
            _add_logo(slide, logo_path)

            if section["tab"] != current_tab:
                current_tab = section["tab"]
                _add_textbox(slide, 0.4, 0.5, 12, 0.6, current_tab, 22, BRAND_DEEP, bold=True)
                y_cursor = 1.3
            else:
                y_cursor = 0.5

            _add_textbox(slide, 0.4, y_cursor, 12.5, 0.6, section["heading"], 13, BRAND_DEEP, bold=True)
            y_charts = y_cursor + 0.7

            figs = section["figs"]
            if not figs:
                continue
            n = len(figs)
            gap = 0.3
            chart_w = (SLIDE_W_IN - 0.8 - gap * (n - 1)) / n
            chart_h = SLIDE_H_IN - y_charts - 0.3
            for j, fig in enumerate(figs):
                img_path = str(tmp_path / f"chart_{i}_{j}.png")
                _fig_to_png(fig, img_path)
                x = 0.4 + j * (chart_w + gap)
                slide.shapes.add_picture(img_path, Inches(x), Inches(y_charts), width=Inches(chart_w), height=Inches(chart_h))

        output = io.BytesIO()
        prs.save(output)
        return output.getvalue()
